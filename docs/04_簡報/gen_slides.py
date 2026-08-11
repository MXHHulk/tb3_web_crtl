# -*- coding: utf-8 -*-
"""
依本專案內容生成一份「從基礎到進階」的教學簡報（PowerPoint .pptx）。

涵蓋範圍（期刊論文等級的詳盡程度）：
  ROS 基礎(節點/話題/pub-sub/action/參數/TF) → 感知與建圖(LiDAR/佔據網格/SLAM/gmapping)
  → 數學基礎(形態學/PCA) → 核心演算法(PCA 對齊牛耕式) → 導航執行(move_base/costmap/DWA)
  → 系統工程(多執行緒/鎖/合作式取消/座標轉換/Flask/Canvas) → 操作/調校/結果/限制/除錯。

輸出：docs/04_簡報/教學簡報.pptx
依 rebuild 分支當前程式碼撰寫。
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.realpath(__file__))

# ───────────────────────────── 配色 ─────────────────────────────
NAVY   = RGBColor(0x0F, 0x2A, 0x43)
TEAL   = RGBColor(0x12, 0x9E, 0x8E)
ORANGE = RGBColor(0xE8, 0x7A, 0x1E)
INK    = RGBColor(0x20, 0x2A, 0x33)
MUTED  = RGBColor(0x6B, 0x7A, 0x88)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0xB8, 0xC6, 0xD0)
CODEBG = RGBColor(0x10, 0x22, 0x33)
CODEFG = RGBColor(0xD6, 0xE6, 0xF2)
GREEN  = RGBColor(0x8B, 0xD4, 0x50)

EA    = "微軟正黑體"      # 中日韓字型
LATIN = "Segoe UI"        # 西文字型
MONO  = "Consolas"        # 等寬字型

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
PAGE = [0]


# ───────────────────────────── 低階工具 ─────────────────────────────
def _ea(run, font=EA):
    """補上 a:ea（東亞字型）屬性，讓中文不會掉回預設字型。"""
    rPr = run._r.get_or_add_rPr()
    el = rPr.find(qn('a:ea'))
    if el is None:
        el = rPr.makeelement(qn('a:ea'), {})
        rPr.append(el)
    el.set('typeface', font)


def _run(p, text, size=18, bold=False, color=INK, latin=LATIN, ea=EA, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = latin
    _ea(r, ea)
    return r


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    return tf


def footer(slide, dark=False):
    col = PALE if dark else MUTED
    tf = textbox(slide, Inches(0.5), Inches(7.04), Inches(11), Inches(0.35))
    _run(tf.paragraphs[0],
         "TurtleBot3 牛耕式完全覆蓋路徑規劃 (CCPP) · 教學簡報", size=9, color=col)
    tf2 = textbox(slide, Inches(12.2), Inches(7.04), Inches(0.9), Inches(0.35))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, str(PAGE[0]), size=9, color=col)


# ─────────────────────── 行內樣式 (**粗體** / `程式`) ───────────────────────
def _emit(p, text, size):
    for seg in re.split(r'(\*\*[^*]+\*\*|`[^`]+`)', text):
        if not seg:
            continue
        if seg.startswith('**') and seg.endswith('**'):
            _run(p, seg[2:-2], size=size, bold=True, color=NAVY)
        elif seg.startswith('`') and seg.endswith('`'):
            _run(p, seg[1:-1], size=size - 1, color=ORANGE, latin=MONO, ea=MONO)
        else:
            _run(p, seg, size=size, color=INK)


# ─────────────────────────── 版面元件 ───────────────────────────
def title_slide():
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    rect(s, 0, Inches(2.05), SW, Inches(0.06), ORANGE)
    rect(s, 0, Inches(5.35), SW, Inches(0.06), TEAL)
    tf = textbox(s, Inches(0.9), Inches(2.3), Inches(11.6), Inches(3.0))
    _run(tf.paragraphs[0], "從基礎到進階的完整教學", size=18, bold=True, color=TEAL)
    p1 = tf.add_paragraph(); p1.space_before = Pt(16)
    _run(p1, "TurtleBot3 牛耕式", size=42, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    _run(p2, "完全覆蓋路徑規劃系統 (CCPP)", size=42, bold=True, color=WHITE)
    p3 = tf.add_paragraph(); p3.space_before = Pt(14)
    _run(p3, "PCA 對齊牛耕式分解 · ROS · SLAM · move_base / DWA · 即時網頁視覺化",
         size=15, color=PALE)
    tf2 = textbox(s, Inches(0.9), Inches(5.6), Inches(11.6), Inches(1.2))
    _run(tf2.paragraphs[0], "MXHHulk", size=16, bold=True, color=WHITE)
    pb = tf2.add_paragraph()
    _run(pb, "機器人與自動化 · ROS Noetic · 依 rebuild 分支當前程式碼撰寫", size=12, color=PALE)
    return s


def section(num, title, subtitle=None):
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    rect(s, Inches(0.9), Inches(3.45), Inches(2.4), Inches(0.05), TEAL)
    tf = textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.0))
    _run(tf.paragraphs[0], f"PART {num}", size=20, bold=True, color=TEAL)
    p1 = tf.add_paragraph(); p1.space_before = Pt(10)
    _run(p1, title, size=40, bold=True, color=WHITE)
    if subtitle:
        p2 = tf.add_paragraph(); p2.space_before = Pt(14)
        _run(p2, subtitle, size=16, color=PALE)
    footer(s, dark=True)
    return s


def content(title, kicker=None):
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(0.12), TEAL)
    tf = textbox(s, Inches(0.6), Inches(0.32), Inches(12.1), Inches(1.15))
    if kicker:
        _run(tf.paragraphs[0], kicker, size=12, bold=True, color=ORANGE)
        pt = tf.add_paragraph()
    else:
        pt = tf.paragraphs[0]
    _run(pt, title, size=26, bold=True, color=NAVY)
    footer(s)
    return s


def body(slide, items, left=0.7, top=1.65, width=11.95, height=5.1, size=17):
    """items: list of (kind, text)。kind ∈ h/b/b2/t/n。text 支援 **粗** 與 `碼`。"""
    tf = textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    first = True
    for kind, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        p.line_spacing = 1.1
        if kind == 'h':
            p.space_before = Pt(8)
            _run(p, text, size=size, bold=True, color=TEAL)
        elif kind == 'b':
            _run(p, "▸  ", size=size, bold=True, color=ORANGE)
            _emit(p, text, size)
        elif kind == 'b2':
            p.level = 1
            _run(p, "–  ", size=size - 2, color=MUTED)
            _emit(p, text, size - 2)
        elif kind == 't':
            _emit(p, text, size)
        elif kind == 'n':
            p.space_before = Pt(4)
            _run(p, text, size=size - 4, italic=True, color=MUTED)
    return tf


def _code_line(p, ln, size):
    p.line_spacing = 1.05
    if '#' in ln:
        i = ln.index('#')
        if ln[:i]:
            _run(p, ln[:i], size=size, color=CODEFG, latin=MONO, ea=MONO)
        _run(p, ln[i:], size=size, color=GREEN, latin=MONO, ea=MONO)
    else:
        _run(p, ln if ln else " ", size=size, color=CODEFG, latin=MONO, ea=MONO)


def code(slide, src, left=0.7, top=1.65, width=11.95, height=None, size=13.5, caption=None):
    lines = src.split('\n')
    if height is None:
        height = 0.30 * len(lines) + 0.45
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.12)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _code_line(p, ln, size)
    if caption:
        ct = textbox(slide, Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.4))
        _run(ct.paragraphs[0], caption, size=11, italic=True, color=MUTED)
    return box


def mono(slide, src, left=0.7, top=1.65, width=11.95, height=None, size=13, color=INK):
    """淺底等寬框，給 ASCII 圖示用。"""
    lines = src.split('\n')
    if height is None:
        height = 0.265 * len(lines) + 0.4
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = PALE; box.line.width = Pt(0.75); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = False; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.1)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.05
        _run(p, ln if ln else " ", size=size, color=color, latin=MONO, ea=MONO)
    return box


def callout(slide, text, top=6.2, color=ORANGE):
    """底部重點提示條。"""
    bar = rect(slide, Inches(0.7), Inches(top), Inches(11.95), Inches(0.72), LIGHT)
    rect(slide, Inches(0.7), Inches(top), Inches(0.1), Inches(0.72), color)
    tf = slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.5), Inches(0.72)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _run(p, "★ ", size=15, bold=True, color=color)
    _emit(p, text, 15)


# ═══════════════════════════════════════════════════════════════════
#  投影片內容
# ═══════════════════════════════════════════════════════════════════
title_slide()

# ── 議程 ──
s = content("議程：九個主題、一條學習路徑")
body(s, [
    ('b', '**PART I — 問題定義**：什麼是完全覆蓋路徑規劃（CCPP），與點對點導航有何不同'),
    ('b', '**PART II — ROS 基礎**：節點、話題、發布/訂閱、Action、參數伺服器、TF 座標轉換'),
    ('b', '**PART III — 感知與建圖**：LiDAR、佔據網格、SLAM、gmapping（粒子濾波）'),
    ('b', '**PART IV — 數學基礎**：形態學（侵蝕/膨脹）、主成分分析（PCA）'),
    ('b', '**PART V — 核心演算法**：PCA 對齊的牛耕式分解與路點生成'),
    ('b', '**PART VI — 導航與執行**：move_base 堆疊、costmap、DWA 動態視窗法'),
    ('b', '**PART VII — 系統工程**：多執行緒、鎖、合作式取消、座標轉換、Flask、Canvas'),
    ('b', '**PART VIII — 操作、調校與結果**：啟動流程、參數調校、結果、限制、除錯'),
    ('n', '每段先講「基礎概念」，再切到「本專案怎麼用、程式在哪」，最後帶「進階細節 / 數學」。'),
], size=16.5)

# ═══════════════ PART I ═══════════════
section("I", "問題定義", "完全覆蓋路徑規劃要解決什麼問題")

s = content("CCPP 是什麼？跟一般導航差在哪", "PART I · 問題定義")
body(s, [
    ('h', '一般導航 (point-to-point)'),
    ('t', '給定起點與終點，求**一條**安全、短的路徑——像 Google Map 導航。'),
    ('h', '完全覆蓋路徑規劃 (Complete Coverage Path Planning, CCPP)'),
    ('t', '要求機器人**走遍工作空間中每一個可到達的點**，理想上不遺漏、不重複。'),
    ('b', '最典型的例子：**掃地機器人**——目標不是「到某點」，而是「把整塊地掃過」。'),
    ('b', '同類問題：地面巡檢、割草、噴漆、農田作業、消毒。'),
    ('b', '兩個品質指標：**覆蓋率**（蓋到多少）與**重複率/總路程**（效率）。'),
], size=17)
callout(s, '本專案 = CCPP + 即時網頁監控 + 一鍵執行，跑在實體 TurtleBot3 Burger 上。')

s = content("覆蓋策略：牛耕式 (Boustrophedon)", "PART I · 問題定義")
body(s, [
    ('t', '名稱取自古希臘「**牛犁田**」：一條線走到底、轉彎、再平行走回來，來回往復。'),
    ('t', '在凸形區域內近乎最佳：路徑簡單、轉彎少、易於實作與分析。'),
], top=1.6, height=1.6, size=17)
mono(s, """  → → → → → → → →
                  ↓
  ← ← ← ← ← ← ← ←
  ↓
  → → → → → → → →      ← 相鄰掃描線間距 = spacing，必須 < 覆蓋寬度才不漏掃""",
     top=3.3, height=2.1, size=14)
callout(s, '核心挑戰：SLAM 地圖的牆面常是「斜的」，固定軸向掃描會變鋸齒狀 → PART IV 用 PCA 解決。')

# ═══════════════ PART II ═══════════════
section("II", "ROS 基礎", "節點、話題、發布/訂閱、Action、參數、TF")

s = content("ROS 是什麼", "PART II · ROS 基礎")
body(s, [
    ('t', 'ROS (Robot Operating System) 不是作業系統，而是一套**讓很多小程式彼此傳訊息**的中介軟體框架。本專案用 **ROS Noetic**。'),
    ('h', '核心哲學：拆成小程式，靠訊息溝通'),
    ('b', '不要把整台機器人寫成一支大程式，而是拆成許多**節點 (Node)**，各做各的事。'),
    ('b', '節點之間透過**話題 (Topic)** 廣播**訊息 (Message)** 來溝通。'),
    ('b', '好處：**解耦、可替換、可重用、可分散在多台電腦**。'),
    ('h', '本專案啟動後同時在跑的節點'),
    ('t', '`gmapping`（建圖）、`move_base`（導航）、`map_server`（規劃+網頁+派工）；`boustrophedon_planner` 為選用、預設停用。'),
], size=16.5)

s = content("節點 / 話題 / 訊息：YouTube 比喻", "PART II · ROS 基礎")
body(s, [
    ('b', '**話題 (Topic)** = YouTube 頻道，有名字（`/map`、`/odom`）與固定的訊息型別。'),
    ('b', '**發布者 (Publisher)** = 上傳影片的人；**訂閱者 (Subscriber)** = 按訂閱的人。'),
    ('b', '發布者與訂閱者**彼此不認識**，只要約好頻道名稱與訊息格式就能溝通 → 這就是「解耦」。'),
    ('h', '本專案用到的話題與訊息型別'),
    ('b2', '`/scan` — `LaserScan`：LiDAR 360° 掃描'),
    ('b2', '`/map` — `OccupancyGrid`：佔據網格地圖（gmapping 發布）'),
    ('b2', '`/odom` — `Odometry`：里程計位置'),
    ('b2', '`/coverage_path` — `nav_msgs/Path`：覆蓋路徑（給 RViz 看，選用節點預設停用）'),
    ('b2', '`/cmd_vel` — `Twist`：速度指令（送給馬達）'),
], size=16)

s = content("Publish：發布一則訊息", "PART II · ROS 基礎")
body(s, [('t', '以 `boustrophedon.py` 為例（此節點預設停用，僅作 publish 範例）：算完覆蓋路徑後發布到 `/coverage_path` 給 RViz。')],
     top=1.55, height=0.8, size=15)
code(s, """# 第 1 步：開一個發布者 (boustrophedon.py:65)
_pub = rospy.Publisher('/coverage_path', Path, queue_size=1, latch=True)
#                       話題名稱          型別   佇列      鎖存(留住最後一則)

# 第 2 步：建訊息、蓋時間戳、填路點，送出去 (boustrophedon.py:53-58)
path = Path()
path.header.stamp    = rospy.Time.now()
path.header.frame_id = msg.header.frame_id or 'map'
path.poses           = [_make_pose(x, y, ...) for x, y in pts]
_pub.publish(path)       # ← 真正送出；所有訂閱者都會收到""", top=2.35, size=13.5)
callout(s, '`latch=True`：晚一點才打開的 RViz 也能立刻收到上一次的路徑，不必空等下一次重算。')

s = content("Subscribe：訂閱 + 回呼 (callback)", "PART II · ROS 基礎")
body(s, [('t', '訂閱 = 跟 ROS 說「這個話題一有新訊息，就自動呼叫我指定的函式」。')],
     top=1.55, height=0.7, size=16)
code(s, """# map_server.py:358-359 — 一次訂閱兩個話題
rospy.Subscriber('/map',  OccupancyGrid, map_callback,  queue_size=1)
rospy.Subscriber('/odom', Odometry,      odom_callback, queue_size=10)
#                 話題     型別            收到要呼叫誰     佇列長度

def map_callback(msg):          # msg 就是收到的那則 OccupancyGrid
    res  = msg.info.resolution  # 從訊息取出解析度
    data = np.array(msg.data).reshape((h, w))   # 取地圖陣列來用

rospy.spin()    # 卡住主執行緒，把 CPU 交給 ROS 去觸發各回呼""", top=2.35, size=13.5)
callout(s, '`/odom` 更新快 → queue_size=10 避免漏；`/map` 又大又慢 → 只留最新 1 則。')

s = content("本專案的 pub/sub 關係圖", "PART II · ROS 基礎")
mono(s, """ 硬體驅動 ──▶ /scan ──▶ gmapping ──▶ /map ──┬──▶ map_server.py (網頁+派工)
 (OpenCR/光達/                                │
  馬達/編碼器/IMU) ──▶ /odom ────────────────┴──▶ map_server.py (記錄軌跡)
        ▲                                          │ 目標點 (Action)
        └───── /cmd_vel ◀── move_base ◀───────────┘

 ┄┄ 選用,預設停用 ┄┄
   /map ──▶ boustrophedon.py ──▶ /coverage_path ──▶ RViz  (功能與網頁重複)""",
     top=1.65, height=2.85, size=12.5)
body(s, [
    ('b', '**只能訂閱「話題」，不能訂閱「節點」**：節點訂的是 `/map`（剛好由 gmapping 發布），不認識也不在乎發布者是誰。'),
    ('b', '一個話題可有**多個訂閱者**：gmapping 發一次 `/map`，所有訂閱者都收到。'),
    ('b', '**RViz 是選用的 ROS 桌面 3D 工具**，與 8080 網頁是兩套；不開 RViz 也能完整使用本專案。'),
], top=4.65, size=15.5)
callout(s, 'boustrophedon_planner 已在 start.launch 預設註解停用——覆蓋路徑改由 map_server.py 計算並顯示於網頁。', top=6.35)

s = content("Action：適合「長時間、可取消」的任務", "PART II · ROS 基礎")
body(s, [
    ('t', '不是所有溝通都走 pub/sub。把目標點交給 `move_base` 用的是 **Action** 機制。'),
    ('h', 'pub/sub vs Action'),
    ('b', '**pub/sub**：持續廣播的資料流（地圖、里程計）——射後不理。'),
    ('b', '**Action**：一次性任務委派，有「**開始 → 過程回饋 → 等結果 → 可中途取消**」。'),
    ('t', '「去 (x, y) 這個點」正好符合：要等它走到、可能失敗、隨時可喊停。'),
], top=1.6, height=3.2, size=16.5)
code(s, """client = actionlib.SimpleActionClient('move_base', MoveBaseAction)
client.send_goal(goal)        # 送目標
client.wait_for_result(...)   # 等結果
client.cancel_all_goals()     # 隨時可取消""", top=5.0, size=13)

s = content("參數伺服器與 TF 座標轉換", "PART II · ROS 基礎")
body(s, [
    ('h', '參數 (Parameter)：節點的設定值'),
    ('b', 'launch 檔用 `<param>` 設定，程式用 `rospy.get_param` 讀取（第二個參數是預設值）。'),
    ('b2', '`port = rospy.get_param(\'~port\', 8080)`  → 網頁埠號，launch 沒給就用 8080。'),
    ('b2', 'launch arg：`model`（burger/waffle）決定載哪一組導航參數檔。'),
    ('h', 'TF：記錄各座標系之間的關係'),
    ('b', '機器人世界裡有多個座標系：`map`（全域）、`odom`（里程計）、`base_link`（車體）。'),
    ('b', 'gmapping 除了發 `/map`，也持續發布 **map→odom 的 TF**，move_base 才知道機器人在地圖上的哪裡。'),
    ('n', 'TF 讓「LiDAR 看到的點」能換算成「地圖上的位置」——SLAM 與導航的黏著劑。'),
], size=16)

# ═══════════════ PART III ═══════════════
section("III", "感知與建圖", "LiDAR、佔據網格、SLAM、gmapping")

s = content("感測器：360° LiDAR", "PART III · 感知與建圖")
body(s, [
    ('b', 'TurtleBot3 頂部的雷射雷達 (LiDAR) 旋轉一圈，量測四周障礙物的**距離與角度**。'),
    ('b', '結果以 `LaserScan` 訊息發布在 `/scan`，是 SLAM 建圖的原始輸入。'),
    ('b', '頻率很重要：太低會建圖品質差。專案附 `test/` 工具量測 `/scan` 頻率。'),
    ('h', '量測 LiDAR 頻率的小工具'),
], top=1.6, height=2.7, size=17)
code(s, """rostopic hz /scan        # 直接看話題頻率
python test/map_server_raw.py   # 專案內的頻率量測腳本""", top=4.3, size=13.5)
callout(s, '感測 → 建圖 → 規劃 → 執行：整條鏈的源頭就是這顆 LiDAR 的掃描資料。')

s = content("佔據網格 OccupancyGrid：地圖的資料結構", "PART III · 感知與建圖")
body(s, [
    ('t', '把空間切成一格一格（像方格紙覆蓋地板），每格存一個整數值：'),
    ('b', '`0` = 自由空間（空地，可走）'),
    ('b', '`100` = 障礙（牆、家具）'),
    ('b', '`-1` = 未知（還沒探索過）'),
    ('h', '訊息裡的關鍵欄位'),
    ('b2', '`info.resolution`：每格幾公尺（典型 0.05 m/格）'),
    ('b2', '`info.width / height`：地圖幾格寬、幾格高'),
    ('b2', '`info.origin`：地圖左下角對應的世界座標'),
    ('b2', '`data`：一維陣列，需 `reshape((h, w))` 還原成二維'),
], top=1.55, size=16.5)
callout(s, '整張地圖 = 一個二維整數陣列。後續所有影像處理、PCA、路徑規劃都建立在它之上。')

s = content("SLAM：邊走邊建圖、同時定位", "PART III · 感知與建圖")
body(s, [
    ('t', 'SLAM = Simultaneous Localization And Mapping（同時定位與建圖）。'),
    ('h', '雞生蛋的兩難'),
    ('b', '要**建地圖**，得先知道自己**在哪**（定位）。'),
    ('b', '要**定位**，又得先有**地圖**當參照。'),
    ('t', 'SLAM 同時解這兩件事：一邊用 LiDAR 掃描更新地圖，一邊用地圖回推自己的位置。'),
    ('h', '本專案用 gmapping'),
    ('b', '輸入 `/scan` + `/odom` → 輸出 `/map`（佔據網格）+ `map→odom` 的 TF。'),
    ('n', '操作上：先「遙控機器人在房間走一圈」，把地圖建得夠完整，再開始覆蓋。'),
], size=16.5)

s = content("gmapping 的原理（進階）", "PART III · 感知與建圖")
body(s, [
    ('h', 'Rao-Blackwellized 粒子濾波器 (RBPF)'),
    ('b', '同時維護**多個「假設」**（粒子），每個粒子代表一條可能的機器人軌跡 + 對應地圖。'),
    ('b', '新掃描進來時，依「掃描和該粒子地圖有多吻合」給每個粒子**加權**。'),
    ('b', '吻合度高的粒子權重高；定期**重採樣**淘汰不準的假設。'),
    ('b', '「Rao-Blackwellized」= 把「軌跡」與「地圖」拆開處理，數學上降低估計變異、提升效率。'),
    ('n', '本專案把 gmapping 當黑盒子使用——它發布 /map，我們訂閱。換成別的 SLAM 只要一樣發 /map，程式一行都不用改（pub/sub 解耦的威力）。'),
], size=16.5)

# ═══════════════ PART IV ═══════════════
section("IV", "數學基礎", "形態學影像處理 · 主成分分析 PCA")

s = content("形態學：侵蝕與膨脹", "PART IV · 數學基礎")
body(s, [
    ('t', '形態學 (Morphology) 是對二值影像「縮小 / 放大」形狀的影像處理運算。'),
    ('h', '兩個基本運算'),
    ('b', '**膨脹 (dilation)**：讓障礙「變胖」一圈 → 用來做安全邊距，路徑遠離牆。'),
    ('b', '**侵蝕 (erosion)**：讓障礙「縮小」 → 用來確認機器人中心可安全到達的範圍。'),
    ('h', '結構元素 (structuring element)'),
    ('b', '決定「胖/瘦多少」的小核，本專案用 `(2r+1)×(2r+1)` 的方形核。'),
    ('b', '半徑（格數） `r = round(margin / resolution)`：把公尺換算成格數。'),
], top=1.55, size=16.5)

s = content("安全邊距：apply_safety_margin()", "PART IV · 數學基礎")
code(s, """def apply_safety_margin(data, margin, resolution):
    obs = data == 100                       # 找出所有障礙格
    r   = max(1, round(margin / resolution))# 邊距(公尺) → 格數
    kern = np.ones((2*r+1, 2*r+1), dtype=bool)  # 方形結構元素
    return binary_dilation(obs, structure=kern) # 障礙膨脹一圈安全距離""", top=1.6, size=13.5)
body(s, [
    ('b', '輸出是 bool 陣列：`True` = 膨脹後的障礙（含安全邊距）。'),
    ('b', '可走集合 = **自由 且 不在膨脹障礙內**：`free = (data == 0) & ~safe_obs`。'),
    ('b', '這一步保證規劃出的路徑**絕不貼牆**，是覆蓋規劃的前處理。'),
], top=3.4, size=16.5)
callout(s, '注意：這是「路徑規劃真正用的」膨脹（MARGIN=0.10 m）；網頁顯示用的三圖層膨脹是另一回事。')

s = content("為什麼需要 PCA？固定軸向會出事", "PART IV · 數學基礎")
body(s, [
    ('t', 'SLAM 地圖的座標方向由「機器人開機那一刻的朝向」決定 → 房間牆面常與 X/Y 軸**斜交**。'),
    ('t', '若沿固定 X 軸來回掃，掃描線與斜牆不平行 → 邊緣變成**破碎的階梯/鋸齒**、轉彎暴增。'),
], top=1.55, height=1.5, size=16.5)
mono(s, """  牆是斜的，掃描線卻是水平的：
     ╱──────────╱
    ╱  → → → →  ╱        每條線長度都不同
   ╱  → → → →  ╱         邊緣全是碎階梯、轉彎多、易漏掃
  ╱──────────╱""", top=3.15, height=1.95, size=13)
callout(s, '解法：用 PCA 找出「房間自己的方向」，讓掃描線自動與牆平行。')

s = content("PCA 的數學（基礎）", "PART IV · 數學基礎")
body(s, [
    ('t', '主成分分析 (Principal Component Analysis)：找出一堆點「散得最開」的方向。'),
    ('h', '三個步驟'),
    ('b', '① 取所有可走格的世界座標，得到點集，計算中心並去中心化。'),
    ('b', '② 算**共變異數矩陣** `C = cov(points)`（2×2，描述點如何分布）。'),
    ('b', '③ 對 C 做**特徵分解** `eigh(C)` → 得到特徵值與特徵向量。'),
    ('h', '幾何意義'),
    ('b', '**最大**特徵值的特徵向量 = **長軸**（點散最開的方向）。'),
    ('b', '**最小**特徵值的特徵向量 = **短軸**（與長軸垂直）。'),
], size=16.5)

s = content("PCA 在程式裡：自動對齊掃描方向", "PART IV · 數學基礎")
code(s, """center           = pts.mean(axis=0)         # 點集中心
diffs            = pts - center             # 去中心化
eigvals, eigvecs = np.linalg.eigh(np.cov(diffs.T))  # 特徵分解
axis_a = eigvecs[:, np.argmax(eigvals)]     # 長軸 → 掃描(sweep)方向
axis_b = eigvecs[:, np.argmin(eigvals)]     # 短軸 → 步進(step)方向""",
     top=1.6, size=13.5, caption="coverage_planner.py:49-54")
body(s, [
    ('b', '沿**長軸**掃描 → 掃描線最長、轉彎最少。'),
    ('b', '每掃完一條，沿**短軸**前進一個 `spacing`。'),
    ('b', '因為軸向是**從資料算出來的**，地圖無論怎麼旋轉，掃描線都自動與牆平行。'),
], top=3.5, size=16.5)
callout(s, '這是本專案最關鍵的設計：用一個經典統計方法，優雅解掉「斜房間」這個工程痛點。')

# ═══════════════ PART V ═══════════════
section("V", "核心演算法", "PCA 對齊的牛耕式分解與路點生成")

s = content("演算法總覽：從地圖到路點", "PART V · 核心演算法")
mono(s, """ /map(佔據網格)
     │  apply_safety_margin()  障礙膨脹 MARGIN
     ▼
  free = (data==0) & ~safe_obs     可走格 (bool 陣列)
     │  boustrophedon()
     ├─ PCA 求長軸/短軸
     ├─ 把可走點投影到長軸(proj_a)、短軸(proj_b)
     ├─ 沿短軸每隔 spacing 取一條掃描線
     └─ 交替方向連成線段
     ▼
  [(wx, wy), ...]  依序走訪的世界座標路點""", top=1.7, height=4.6, size=13)
callout(s, '純演算法在 coverage_planner.py，完全不依賴 ROS → 可獨立測試、被多個節點共用。')

s = content("掃描線生成與方向交替", "PART V · 核心演算法")
code(s, """b = b_min + spacing / 2          # 第一條線從邊界內縮半格距,兩端對稱
while b <= b_max + spacing / 2:
    mask = np.abs(proj_b - b) < spacing / 2   # 落在這條線附近的點
    if mask.any():
        a_vals  = proj_a[mask]
        p_start = center + a_vals.min() * axis_a + b * axis_b
        p_end   = center + a_vals.max() * axis_a + b * axis_b
        if l2r:  waypoints += [p_start, p_end]   # 左→右
        else:    waypoints += [p_end, p_start]   # 右→左 (反向)
        l2r = not l2r                # 下一條反過來 → 牛耕來回
    b += spacing""", top=1.6, size=12.5, caption="coverage_planner.py:67-82")
body(s, [
    ('b', '每條掃描線取「最遠的兩個投影端點」連成線段。'),
    ('b', '布林旗標 `l2r` 讓相鄰線方向**交替**——這正是牛耕來回的精髓。'),
], top=5.45, size=16)

s = content("關鍵參數：SPACING 與 MARGIN", "PART V · 核心演算法")
code(s, """SPACING = 0.18   # 掃描線間距(公尺);覆蓋寬度 ≈ 0.20 m,留 10% 重疊
MARGIN  = 0.10   # 障礙物安全邊距(公尺) ≈ 機器人半徑(直徑 0.20 m)""",
     top=1.55, size=14, caption="coverage_planner.py:8-9（rebuild 分支當前值）")
body(s, [
    ('h', 'SPACING（掃描線間距）'),
    ('b', '必須 **< 機器人覆蓋寬度** 才會重疊、不漏掃。實測覆蓋寬度 ≈ 0.20 m。'),
    ('b', '設 0.18 m → 相鄰線約留 10% 重疊。調小：覆蓋率↑、但掃描線變多、變慢。'),
    ('h', 'MARGIN（安全邊距）'),
    ('b', '≈ 機器人半徑（直徑 0.20 m → 半徑 0.10 m）。調大：更安全但邊角漏掃；調小：更貼牆但風險高。'),
    ('n', '經驗：MARGIN=0.10 已是「邊緣剛好貼牆」的最小值；實機若擦牆或規劃常失敗，回加到 0.11–0.12。'),
], top=2.55, size=16)

s = content("案例：間距過大如何造成覆蓋率偏低", "PART V · 核心演算法")
body(s, [
    ('h', '問題現象'),
    ('b', '原設定 `SPACING = 0.25 m`，但量測覆蓋寬度只有 0.20 m。'),
    ('b', '0.25 > 0.20 → 相鄰掃描線間留下約 **0.05 m 的未覆蓋帶**，覆蓋率看起來偏低。'),
    ('h', '修正'),
    ('b', '`SPACING` 0.25 → **0.18 m**：掃描線變密，改為約 0.02 m 重疊，不再留縫。'),
    ('b', '`MARGIN` 0.15 → **0.10 m**：邊距從「比半徑大 50%」收到貼近實際半徑，靠牆區能覆蓋。'),
    ('n', '兩個常數都在 coverage_planner 模組，boustrophedon.py 與 map_server.py 都 import 它 → 改一處全生效。'),
], size=16.5)

# ═══════════════ PART VI ═══════════════
section("VI", "導航與執行", "move_base 堆疊 · costmap · DWA")

s = content("move_base：導航的總管", "PART VI · 導航與執行")
body(s, [
    ('t', 'move_base 是 ROS 的**導航堆疊**。你給它一個目標點，它負責「把機器人安全開過去」整件事。'),
    ('h', '內部把四塊串起來'),
    ('b', '**Global planner**：看整張地圖，畫一條大致路線（像先規劃整段路）。'),
    ('b', '**Local planner**（就是 DWA）：盯著眼前幾秒，即時算速度、閃突發障礙。'),
    ('b', '**Costmap**（全域 + 局部）：障礙 + 膨脹層，分別給上面兩個規劃器用。'),
    ('b', '**Recovery behaviors**：卡住時的脫困（原地旋轉、清 costmap）。'),
    ('n', 'start.launch 啟動「一個」move_base 節點，再塞入 turtlebot3_navigation 的標準參數檔。'),
], size=16.5)

s = content("costmap：成本地圖與膨脹層", "PART VI · 導航與執行")
body(s, [
    ('b', 'costmap 在佔據網格之上，替每格標一個「**通過成本**」：越靠近障礙成本越高。'),
    ('b', '**膨脹層 (inflation)**：障礙周圍漸層加成本，讓規劃器自然避開牆邊。'),
    ('h', '兩種 costmap'),
    ('b', '**Global costmap**：整張地圖（吃 `/map`，static_map），給全域規劃器。'),
    ('b', '**Local costmap**：跟著機器人移動的小窗（吃 LiDAR），給 DWA 即時避障。'),
    ('n', '本專案不自訂這些參數，直接沿用 turtlebot3_navigation 套件的 YAML（依 model 自動選）。'),
], top=1.6, size=16.5)
callout(s, '想調「移動速度 / 到點容差」要改 turtlebot3_navigation 的 YAML，不是本專案的檔。')

s = content("DWA：動態視窗法（進階）", "PART VI · 導航與執行")
body(s, [
    ('t', 'DWA = Dynamic Window Approach，move_base 內負責**即時避障**的局部規劃器。'),
    ('h', '運作流程（每幾十毫秒一輪）'),
    ('b', '① 在「當前速度 + 加速度上限」可達範圍（=**動態視窗**）內，取樣多組 (v, ω)。'),
    ('b', '② 模擬每組「這樣走一小段會到哪」。'),
    ('b', '③ 用評分函式打分：**離障礙夠遠？朝全域路線前進？速度夠快？**'),
    ('b', '④ 選最高分那組 → 輸出成 `/cmd_vel` 發給馬達。'),
    ('n', 'dwa_local_planner_params_<model>.yaml 裡的速度/加速度上限，就是在限制第①步那個視窗多大。'),
], size=16.5)

s = content("派工：把路點送給 move_base", "PART VI · 導航與執行")
body(s, [('t', '`map_server.py` 扮演「派工的人」，只跟 move_base 打交道，碰不到也不需要懂 DWA。')],
     top=1.55, height=0.7, size=16)
code(s, """client = actionlib.SimpleActionClient('move_base', MoveBaseAction)
# ...每個路點:
goal = MoveBaseGoal()
goal.target_pose.header.frame_id = fid
goal.target_pose.pose.position.x = x        # 目標座標
goal.target_pose.pose.position.y = y
goal.target_pose.pose.orientation.z = qz    # 朝向下一個路點
client.send_goal(goal)                       # 「這個點,你想辦法開過去」""",
     top=2.35, size=13, caption="map_server.py:182, 212-219")
callout(s, '分工：map_server 說「去哪、到了沒」；move_base/DWA 管「怎麼走、怎麼閃、馬達轉多快」。')

s = content("失敗容錯：單點到不了不中斷", "PART VI · 導航與執行")
code(s, """if client.get_state() != GoalStatus.SUCCEEDED:
    rospy.logwarn(f'[coverage] 路點 {i+1}/{n} 跳過(狀態碼 {...})')
    cov_status['msg'] = f'路點 {i+1}/{n} 跳過(無法到達)'
    # 不 return,不 raise → 繼續下一點""", top=1.6, size=13.5,
     caption="map_server.py:231-234")
body(s, [
    ('b', '某路點 move_base 回傳非 `SUCCEEDED`（定位漂移、暫時障礙）時，**記錄警告後跳過**。'),
    ('b', '單一無法到達的目標，不會讓整個覆蓋任務卡死或中止。'),
    ('b', '只需看 move_base 回傳的狀態碼 → 完全不用管底層避障細節。'),
], top=3.5, size=16.5)

# ═══════════════ PART VII ═══════════════
section("VII", "系統工程", "多執行緒 · 鎖 · 座標轉換 · Flask · Canvas")

s = content("系統總覽：資料如何流動", "PART VII · 系統工程")
mono(s, """ 實體 TurtleBot3 (LiDAR /scan + 馬達 + 編碼器/IMU)
        │ /scan                          │ /odom
        ▼                                ▼
   gmapping (SLAM) ──▶ /map + TF ──▶ map_server.py
                                     (網頁圖層 + 記錄軌跡 + 派工)
                                          │ 呼叫 coverage_planner 算路點
                                          │ 目標點 (Action)
                                          ▼
                          move_base ──▶ DWA ──▶ /cmd_vel ──▶ 馬達

 ┄ 選用,預設停用 ┄  /map ──▶ boustrophedon.py ──▶ /coverage_path ──▶ RViz
                    (同樣呼叫 coverage_planner,僅供 RViz 顯示)""",
     top=1.65, height=4.9, size=12)
callout(s, 'LiDAR→gmapping 建圖→規劃成路點→map_server 逐點派給 move_base→避障走過去→網頁即時顯示。')

s = content("模組拆解：演算法與 ROS 解耦", "PART VII · 系統工程")
body(s, [
    ('h', 'coverage_planner.py — 純演算法核心 ★'),
    ('b', '只用 numpy/scipy，輸入輸出都是普通陣列與座標，**完全不依賴 ROS**。'),
    ('b', '好處：可獨立單元測試、可被 `boustrophedon.py` 與 `map_server.py` 共用。'),
    ('h', 'boustrophedon.py — 選用規劃節點 (Module B，預設停用)'),
    ('b', '訂 `/map`、每 5 秒重算、發 `/coverage_path` 給 RViz。功能與網頁重複，已在 start.launch 註解停用。'),
    ('h', 'map_server.py — 地圖伺服器 + 派工 (Module A，核心樞紐)'),
    ('b', '地圖轉三圖層 PNG、記錄軌跡、Flask 網頁、把路點派給 move_base。'),
    ('h', 'web/index.html — 前端；start.launch — 一鍵啟動全系統'),
], size=16)

s = content("併發：三條執行緒、三把鎖", "PART VII · 系統工程")
body(s, [
    ('t', '`map_server.py` 同時跑三類執行緒，共享資料必須防資料競爭。'),
    ('b', '**ROS 回呼緒**：`map_callback` / `odom_callback` 不斷更新地圖與位置。'),
    ('b', '**Flask 緒**：處理網頁請求時要讀共享資料。'),
    ('b', '**任務緒**：`run_coverage` 在背景把路點送給 move_base。'),
], top=1.55, height=2.2, size=16.5)
code(s, """map_lock   = threading.Lock()   # 保護地圖資料
robot_lock = threading.Lock()   # 保護機器人位置/軌跡
cov_lock   = threading.Lock()   # 保護覆蓋任務狀態
# 每次讀寫共享資料都用  with xxx_lock:  包起來""", top=3.9, size=13.5)
callout(s, '用鎖確保 Flask 不會讀到「寫到一半」的地圖或狀態。')

s = content("合作式取消 (cooperative cancellation)", "PART VII · 系統工程")
body(s, [('t', '停止任務不是「強制 kill 執行緒」，而是迴圈裡定期檢查狀態、主動優雅退出。')],
     top=1.55, height=0.7, size=16)
code(s, """while not rospy.is_shutdown():
    with cov_lock:
        if cov_status['state'] != 'running':
            client.cancel_all_goals()   # 取消 move_base 目標
            return                       # 優雅退出
    if client.wait_for_result(rospy.Duration(0.5)):
        break                            # 每 0.5 秒檢查一次停止訊號""",
     top=2.35, size=13.5, caption="map_server.py:222-228")
callout(s, '好處：資源乾淨釋放、move_base 目標被正確取消，不會留下殭屍任務。')

s = content("座標轉換：世界 → 格 → 像素", "PART VII · 系統工程")
body(s, [('t', '三種座標：世界(公尺) / 網格(格索引) / 圖片像素。網頁要把世界座標畫到 Canvas。')],
     top=1.55, height=0.7, size=16)
code(s, """px = (wx - origin_x) / resolution - c0
py = crop_h - ((wy - origin_y) / resolution - r0)   # Y 軸翻轉""",
     top=2.35, size=13.5, caption="map_server.py:244-245")
body(s, [
    ('b', '`(世界座標 − 原點) / 解析度` = 格座標；再扣裁切左/上邊界。'),
    ('b', '**Y 軸要翻轉**：影像原點在左上、Y 向下；世界 Y 向上（`np.flipud` 同理）。'),
    ('b', '動態裁切 `_crop`：只保留「已知區」外接矩形 + `CROP_PAD=10` 格，減少傳輸量。'),
], top=3.4, size=16.5)

s = content("Flask 網頁伺服器 + 前端 Canvas", "PART VII · 系統工程")
body(s, [
    ('h', '後端：Flask REST API（map_server.py）'),
    ('b2', '`/map.png` `/map_eroded.png` `/map_dilated.png`：三種地圖圖層'),
    ('b2', '`/robot_state`：機器人位置 + 軌跡（像素座標）'),
    ('b2', '`/coverage/start | stop | status`：一鍵控制與進度查詢'),
    ('h', '前端：HTML5 Canvas 多圖層疊加（index.html）'),
    ('b', '原始地圖打底 → 侵蝕/膨脹層半透明著色疊上 → 再畫覆蓋路徑、行走軌跡、機器人位置。'),
    ('b', '每秒輪詢一次後端（`setInterval(refresh, 1000)`）刷新；圖層可一鍵開關。'),
    ('n', '三圖層：白=空地、黑=障礙、灰=未知；侵蝕(障礙縮小)/膨脹(障礙放大)僅供人眼判讀。'),
], size=15.5)

# ═══════════════ PART VIII ═══════════════
section("VIII", "操作、調校與結果", "啟動流程 · 參數調校 · 結果 · 限制 · 除錯")

s = content("啟動與操作流程", "PART VIII · 操作")
code(s, """pip3 install flask Pillow            # 前置(numpy 已隨 ROS)
roslaunch turtlebot3_ccpp start.launch              # 預設 burger
roslaunch turtlebot3_ccpp start.launch model:=waffle # 換型號
roslaunch turtlebot3_ccpp start.launch port:=9000    # 換埠號""",
     top=1.55, size=13)
body(s, [
    ('b', '① 啟動後先**遙控機器人在房間走一圈**，讓 gmapping 把地圖建完整。'),
    ('b', '② 瀏覽器開 `http://<機器人IP>:8080`，看地圖、位置、規劃出的覆蓋路徑。'),
    ('b', '③ 點 **▶ 開始覆蓋**，機器人依路點自動走遍全場；隨時可點 **■ 停止**。'),
], top=3.3, size=16.5)

s = content("參數調校指南：覆蓋率 vs 速度", "PART VIII · 操作")
body(s, [
    ('h', '想提高覆蓋率（漏掃）'),
    ('b', '調小 `SPACING`（如 0.18→0.15）：掃描線更密、重疊更多，但變慢。'),
    ('b', '調小 `MARGIN`（如 0.12→0.10）：路徑更貼牆，邊角覆蓋↑（但別小於半徑）。'),
    ('h', '想加快 / 減少重複'),
    ('b', '調大 `SPACING`：但不可超過覆蓋寬度，否則重新漏掃。'),
    ('h', '撞牆 / 規劃常失敗'),
    ('b', '調大 `MARGIN`，並檢查 TF 與定位是否漂移。'),
    ('n', '所有覆蓋參數集中在 coverage_planner.py 頂部；導航速度在 turtlebot3_navigation 的 YAML。'),
], size=16)

s = content("結果與討論", "PART VIII · 結果")
body(s, [
    ('h', '實驗（實體 TurtleBot3 Burger，ROS Noetic）'),
    ('b', '建圖完成後，規劃出的掃描線**穩定地與房間牆面平行**，即使地圖座標相對房間旋轉亦然。'),
    ('b', '證實 **PCA 對齊**消除了固定軸向掃描的階梯鋸齒，減少轉彎、提升覆蓋均勻度。'),
    ('b', '安全邊距膨脹使執行軌跡**遠離牆面**；逐點容錯讓個別目標暫時到不了時仍能完成覆蓋。'),
    ('h', '工程上的價值'),
    ('b', '演算法與 ROS 解耦 → 覆蓋邏輯可**獨立於機器人**演練、調校與測試。'),
], size=16.5)

s = content("限制與未來工作", "PART VIII · 結果")
body(s, [
    ('h', '目前限制'),
    ('b', '仰賴**單一連通自由區**：PCA 對齊對「具主導性的矩形房間」最有效，非凸/多房間較弱。'),
    ('b', '固定 **S 形路點順序**，未就總行程做全域最佳化（接近 TSP 問題）。'),
    ('h', '未來方向'),
    ('b', '非凸環境的**多區域分解**（完整臨界點 boustrophedon decomposition）。'),
    ('b', '路點順序的**行程最佳化**（TSP 形式）。'),
    ('b', '**frontier 自主探索**：省去人工建圖步驟。'),
    ('b', '網頁上**即時覆蓋率指標**（實走格 / 應覆蓋格）。'),
], size=16.5)

s = content("常見問題與除錯", "PART VIII · 除錯")
body(s, [
    ('b', '網頁顯示「等待地圖」→ `/map` 未發布；查 gmapping、`rostopic echo /map`。'),
    ('b', '「無可走路徑」→ 地圖太空或全障礙；先多遙控建圖。'),
    ('b', '「move_base 未啟動」→ 看 launch 輸出、`rosnode list`。'),
    ('b', '撞牆 → 調大 `MARGIN`；漏掃 → 調小 `SPACING`。'),
    ('b', '路點一直跳過 → 邊距太大堵死路、或定位漂移；調 `MARGIN`、檢查 TF。'),
    ('h', '好用指令'),
    ('t', '`rostopic list` · `rostopic hz /scan` · `rosnode list` · `rqt_graph`'),
], top=1.6, size=16)

# ═══════════════ 結語 ═══════════════
section("結語", "重點回顧與參考文獻")

s = content("一頁總結：四個設計亮點", "結語")
body(s, [
    ('b', '**① 演算法與 ROS 解耦**：純 numpy/scipy 核心，可測試、可重用。'),
    ('b', '**② PCA 自動對齊**：用經典統計方法，優雅解掉「斜房間」的工程痛點。'),
    ('b', '**③ 多執行緒安全 + 合作式取消**：三把鎖護共享狀態，停止任務乾淨優雅。'),
    ('b', '**④ 多圖層網頁視覺化**：Canvas 疊加地圖/路徑/軌跡/位置，一鍵控制。'),
    ('h', '一句話流程'),
    ('t', 'LiDAR 掃描 → gmapping 建圖 → PCA+牛耕 規劃路點 → move_base/DWA 避障執行 → 網頁即時監控。'),
    ('n', '參數重點：覆蓋寬度 ≈ 0.20 m → SPACING 0.18（留 10% 重疊）、MARGIN 0.10（≈ 半徑）。'),
], size=16.5)

s = content("參考文獻", "結語")
refs = [
    'E. Galceran, M. Carreras, "A survey on coverage path planning for robotics," Robotics and Autonomous Systems, 61(12):1258–1276, 2013.',
    'H. Choset, P. Pignon, "Coverage path planning: The boustrophedon cellular decomposition," Field and Service Robotics, Springer, 1998.',
    'G. Grisetti, C. Stachniss, W. Burgard, "Improved techniques for grid mapping with Rao-Blackwellized particle filters," IEEE T-RO, 23(1):34–46, 2007.',
    'D. Fox, W. Burgard, S. Thrun, "The dynamic window approach to collision avoidance," IEEE Robotics & Automation Magazine, 4(1):23–33, 1997.',
    'M. Quigley et al., "ROS: an open-source Robot Operating System," ICRA Workshop on Open Source Software, 2009.',
    'R. Bormann et al., "Indoor coverage path planning: Survey, implementation, analysis," IEEE ICRA, 2018, pp. 1718–1725.',
    'I. T. Jolliffe, J. Cadima, "Principal component analysis: a review and recent developments," Phil. Trans. R. Soc. A, 374(2065), 2016.',
]
items = []
for i, r in enumerate(refs, 1):
    items.append(('t', f'[{i}]  {r}'))
body(s, items, top=1.6, size=13.5)

s = prs.slides.add_slide(BLANK); bg(s, NAVY); PAGE[0] += 1
rect(s, Inches(0.9), Inches(3.7), Inches(2.4), Inches(0.05), TEAL)
tf = textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.0))
_run(tf.paragraphs[0], "謝謝聆聽 · Q & A", size=40, bold=True, color=WHITE)
pb = tf.add_paragraph(); pb.space_before = Pt(16)
_run(pb, "TurtleBot3 牛耕式完全覆蓋路徑規劃 (CCPP)  ·  MXHHulk", size=15, color=PALE)

# ───────────────────────────── 輸出 ─────────────────────────────
out = os.path.join(HERE, "教學簡報.pptx")
prs.save(out)
print(f"生成完成（共 {PAGE[0]} 張投影片）：")
print(" ", out)
